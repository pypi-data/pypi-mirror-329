"""Extract text and images from PPTX files."""

import concurrent.futures
import io
import os
from dataclasses import dataclass
from typing import List, Tuple

from PIL import Image
from pptx import Presentation

from pyvisionai.extractors.base import BaseExtractor
from pyvisionai.utils.logger import logger


@dataclass
class ImageTask:
    """Container for image processing task data."""

    image_data: bytes
    image_name: str
    output_dir: str
    index: int


class PptxTextImageExtractor(BaseExtractor):
    """Extract text and images from PPTX files."""

    def extract_text_and_images(
        self, pptx_path: str
    ) -> Tuple[List[str], List[bytes]]:
        """Extract text and images from PPTX file."""
        prs = Presentation(pptx_path)
        texts = []
        images = []

        # Extract text and images from slides
        for slide in prs.slides:
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            texts.append("\n".join(slide_text))

            # Extract images from relationships
            for rel in slide.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        images.append(image_data)
                    except Exception as e:
                        logger.error(
                            f"Error extracting image: {str(e)}"
                        )
                        continue

        return texts, images

    def save_image(
        self, image_data: bytes, output_dir: str, image_name: str
    ) -> str:
        """Save an image to the output directory."""
        try:
            # Convert image data to PIL Image
            image = Image.open(io.BytesIO(image_data))
            # Convert to RGB if necessary
            if image.mode != "RGB":
                image = image.convert("RGB")
            # Save as JPEG (supported format)
            img_path = os.path.join(output_dir, f"{image_name}.jpg")
            image.save(img_path, "JPEG", quality=95)
            return img_path
        except Exception as e:
            logger.error(f"Error saving image: {str(e)}")
            raise

    def process_image_task(self, task: ImageTask) -> tuple[int, str]:
        """Process a single image task."""
        try:
            img_path = self.save_image(
                task.image_data, task.output_dir, task.image_name
            )
            image_description = self.describe_image(img_path)
            os.remove(img_path)  # Clean up
            return task.index, image_description
        except Exception as e:
            logger.error(
                f"Error processing image {task.image_name}: {str(e)}"
            )
            return (
                task.index,
                f"Error: Could not process image {task.image_name}",
            )

    def extract(self, pptx_path: str, output_dir: str) -> str:
        """Process PPTX file by extracting text and images separately."""
        try:
            pptx_filename = os.path.splitext(
                os.path.basename(pptx_path)
            )[0]

            # Extract text and images
            texts, images = self.extract_text_and_images(pptx_path)

            # Generate markdown content
            md_content = f"# {pptx_filename}\n\n"

            # Add text content
            for slide_num, text in enumerate(texts, 1):
                if text:
                    md_content += f"## Slide {slide_num}\n\n"
                    md_content += f"{text}\n\n"

            # Prepare image tasks
            image_tasks = []
            for img_index, img_data in enumerate(images):
                image_name = f"{pptx_filename}_image_{img_index + 1}"
                task = ImageTask(
                    image_data=img_data,
                    image_name=image_name,
                    output_dir=output_dir,
                    index=img_index,
                )
                image_tasks.append(task)

            # Process images in parallel if there are any
            if image_tasks:
                # Store descriptions in order
                descriptions = [""] * len(image_tasks)

                # Use ThreadPoolExecutor for parallel processing
                with concurrent.futures.ThreadPoolExecutor(
                    max_workers=4
                ) as executor:
                    # Submit all tasks
                    future_to_task = {
                        executor.submit(
                            self.process_image_task, task
                        ): task
                        for task in image_tasks
                    }

                    # Collect results as they complete
                    for future in concurrent.futures.as_completed(
                        future_to_task
                    ):
                        idx, description = future.result()
                        descriptions[idx] = description

                # Add descriptions to markdown in correct order
                for img_index, description in enumerate(descriptions):
                    md_content += f"[Image {img_index + 1}]\n"
                    md_content += f"Description: {description}\n\n"

            # Save markdown file
            md_file_path = os.path.join(
                output_dir, f"{pptx_filename}_pptx.md"
            )
            with open(md_file_path, "w", encoding="utf-8") as md_file:
                md_file.write(md_content)

            # Add info logging
            logger.info("Processing PPTX file...")
            logger.info(f"Extracted {len(images)} images")
            logger.info("PPTX processing completed successfully")

            return md_file_path

        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise
