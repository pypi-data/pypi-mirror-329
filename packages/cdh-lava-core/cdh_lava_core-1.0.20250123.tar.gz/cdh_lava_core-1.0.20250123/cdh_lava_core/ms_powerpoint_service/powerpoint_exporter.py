from pptx import Presentation
import os
import sys
from cdh_lava_core.cdc_log_service.environment_logging import LoggerSingleton

from cdh_lava_core.cdc_tech_environment_service import environment_file as cdc_env_file
from cdh_lava_core.cdc_tech_environment_service.environment_http import EnvironmentHttp

# Get the currently running file name
NAMESPACE_NAME = os.path.basename(os.path.dirname(__file__))
# Get the parent folder name of the running file
SERVICE_NAME = os.path.basename(__file__)

class PowerPointExporter:
    """
    A class that provides functionality to download and export PowerPoint slides to PNG files.
    """

    @staticmethod
    def save_slide(slide, output_file, data_product_id, environment):
        """
        Function to save a slide as a PNG file.

        Args:
            slide: The slide object to be saved.
            output_file: The path and filename of the output PNG file.

        Returns:
            None
        """
        
            tracer, logger = LoggerSingleton.instance(
                NAMESPACE_NAME, SERVICE_NAME, data_product_id, environment
            ).initialize_logging_and_tracing()

            with tracer.start_as_current_span("get_tasks"):
                try:
                    
                    # Implement function to save a slide as PNG
                    pass

                except Exception as ex_:
                    error_msg = "Error: %s", str(ex_)
                    exc_info = sys.exc_info()
                    LoggerSingleton.instance(
                        NAMESPACE_NAME, SERVICE_NAME, data_product_id, environment
                    ).error_with_exception(error_msg, exc_info)
                    raise


    @staticmethod
    def get_slide_title(slide):
        """
        Function to get the title of a slide.

        Parameters:
        slide (Slide): The slide object from which to extract the title.

        Returns:
        str or None: The title of the slide if found, None otherwise.
        """
        # Search for the title placeholder and return its text
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                return shape.text
        return None

    @classmethod
    def export_slides(cls, pptx_file, output_folder):
        """Function to export all slides to PNG files."""
        prs = Presentation(pptx_file)
        for i, slide in enumerate(prs.slides):
            title = cls.get_slide_title(slide) or f"Slide{i+1}"
            output_file = os.path.join(output_folder, f"{title}.png")
            cls.save_slide(slide, output_file)

    if __name__ == "__main__":
        pptx_file = "path/to/your/presentation.pptx"
        output_folder = "path/to/output/folder"
        export_slides(pptx_file, output_folder)
