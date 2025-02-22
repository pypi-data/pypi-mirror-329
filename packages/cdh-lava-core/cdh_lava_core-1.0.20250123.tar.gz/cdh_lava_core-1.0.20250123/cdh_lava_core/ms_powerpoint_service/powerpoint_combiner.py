from pptx import Presentation
import os
import sys
import io
from cdh_lava_core.cdc_log_service.environment_logging import LoggerSingleton
from cdh_lava_core.cdc_tech_environment_service import environment_file as cdc_env_file
from cdh_lava_core.cdc_tech_environment_service.environment_http import EnvironmentHttp

# Get the currently running file name
NAMESPACE_NAME = os.path.basename(os.path.dirname(__file__))
# Get the parent folder name of the running file
SERVICE_NAME = os.path.basename(__file__)

class PowerPointCombiner:
    @staticmethod
    def combine_presentations(directory_path, output_filename, data_product_id, environment):
        # Create a new presentation
        merged_pptx = Presentation()
        
        # Get all PowerPoint files in directory and sort them
        pptx_files = sorted(
            (f for f in os.listdir(directory_path) 
            if f.endswith(('.pptx', '.ppt')) 
            and f != 'merged.pptx'
            and not f.startswith('~')
            ), 
            key=str.lower
        )
        
        # Iterate through each presentation
        for pptx_file in pptx_files:
            try:
                file_path = os.path.join(directory_path, pptx_file)
                prs = Presentation(file_path)
                
                # Copy each slide from the current presentation
                for slide in prs.slides:
                    # Create a blank slide in the merged presentation with the same layout
                    merged_slide = merged_pptx.slides.add_slide(slide.slide_layout)
                    
                    # Copy the background if it exists
                    # if slide.background:
                    #     merged_slide.background = slide.background
                    
                    # Copy all shapes and their properties, including pictures
                    for shape in slide.shapes:
                        pic = None
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            image_stream = io.BytesIO(shape.image.blob)
                            # Get the original image information
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            # Add the image to the new slide
                            pic = merged_slide.shapes.add_picture(
                                image_stream,
                                left,
                                top,
                                width,
                                height
                            )
                        else:
                            el = shape.element
                            merged_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
            
            except Exception as e:
                exc_info = sys.exc_info()
                LoggerSingleton.instance(
                    SERVICE_NAME, NAMESPACE_NAME, data_product_id, environment
                ).error_with_exception(
                    f"Error processing file {pptx_file}: {str(e)}",
                    exc_info
                )
                raise
        
        try:
            # Save the merged presentation
            merged_pptx.save(output_filename)
        except Exception as e:
            exc_info = sys.exc_info()
            LoggerSingleton.instance(
                    SERVICE_NAME, NAMESPACE_NAME, data_product_id, environment
                ).error_with_exception(
                f"Error saving merged presentation: {str(e)}",
                exc_info
            )
            raise

    @staticmethod
    def get_pptx_directory() -> str:
        """
        Gets the path to the pptx directory relative to the project root.
        
        Returns:
            str: Absolute path to the pptx directory
        """
        # Get the current file's directory
        current_dir = os.path.dirname(os.path.abspath(__file__))
        
        # Navigate up to the project root (parent of tests directory)
        project_root = os.path.dirname(os.path.dirname(current_dir))
        
        # Construct path to pptx directory
        pptx_dir = os.path.join(project_root, 'pptx')
        
        # Create directory if it doesn't exist
        if not os.path.exists(pptx_dir):
            os.makedirs(pptx_dir)
            
        return pptx_dir

    @staticmethod
    def combine_project_presentations(data_product_id, environment, output_filename: str = None) -> str:
        """
        Combines all PowerPoint files in the project's pptx directory.
        
        Args:
            output_filename (str, optional): Name of the output file. If not provided,
                                          defaults to 'merged.pptx' in the pptx directory.
        
        Returns:
            str: Path to the merged presentation file
        """
        pptx_dir = PowerPointCombiner.get_pptx_directory()
        
        if output_filename is None:
            output_filename = os.path.join(pptx_dir, 'merged.pptx')
            
        PowerPointCombiner.combine_presentations(pptx_dir, output_filename, data_product_id, environment)
        return output_filename