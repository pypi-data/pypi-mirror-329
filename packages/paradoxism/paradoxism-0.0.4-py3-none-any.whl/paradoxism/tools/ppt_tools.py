import json
import os
from paradoxism.context import *
from paradoxism.tools import *
from paradoxism.ops.base import prompt
from paradoxism.utils import make_dir_if_need,split_path,yellow_color
from paradoxism.tools.image_tools import text2im
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image

cxt=_context()


@tool('gpt-4o')
def generate_ppt_outlines(topic:str,pages:int=20, main_style='專業的包浩斯風格',language='zh-TW'):
    ppt_prompt=f"""你是一個簡報大師，你擅長將複雜的題目以最吸睛的方式表達，請基於第一性原理思考{topic}作為簡報題目的根本目的，然後請上網查詢這主題，根據查詢後結果以及你對題目的理解，藉此設計出以{language}撰寫，總共{pages}頁的投影片大綱，包含封面頁，最後的Q&A頁，不需要目錄頁。大綱每頁需要:
    - 頁碼(page_num)
    - 主標題(title)
    - 副標題(subtitle, optional)
    - 內文(content)，也就是要顯示在投影片上的文字，應該要能完整的表達知識點或觀念，可透過無序列表的形式來呈現重點，切忌文字量過多
    - 背景圖(background image,optional，如果要請基於{main_style}撰寫出這張圖的生成prompt，如果是有背景圖，應該該畫面只會有主標題或是副標題，且以白色粗體大型字墊在背景圖前,有背景圖時就不該使用插圖)
    - 輔助視覺化，在投影片中可能有以下幾種輔助視覺化技術(整份投影片規劃，每種輔助視覺化技術至少都該出現一次)，只需要列出有出現的:
        - 插圖(image)，請撰寫出基於{main_style}該插圖生成的prompt，避免出現文字，也要考慮插圖的意義是否與該頁內容匹配，且不同頁面插圖間應該要有一定的差異，同時還需要加入size, position兩個屬性，size是指預計要將圖片縮放成哪個尺寸， position則是圖片的左上角所在的位置座標
        - 表格(table)，請以markdown的形式來表達此表格，需註明出處來源
        - 圖表(graph)，請以matplotlib來繪製圖表，請撰寫出產生此圖表的python語法，圖表內描述文字應與使用者常用語言相同
    以json陣列的形式輸出，直接輸出，無須解釋
    """
    ppt_json=prompt(ppt_prompt)

    with open(os.path.join(cxt.get_paradoxism_dir(),'ppt_outlines_{0}.json'.format(get_time_suffix())),'w') as ff:
        ff.write(json.dumps(ppt_json,ensure_ascii=False,indent=4))
    print(yellow_color(ppt_json),flush=True)
    return ppt_json,os.path.join(cxt.get_paradoxism_dir(),'ppt_outlines_{0}.json'.format(get_time_suffix()))



def initialize_presentation(title, subtitle):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 長寬比，寬度
    prs.slide_height = Inches(7.5)   # 16:9 長寬比，高度

    # 添加封面頁
    cover_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(cover_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle



    return prs


# 2. 建議基礎頁面
def suggest_base_slide(prs, page_number, layout_type, title, content):
    if layout_type == "text_only":
        slide_layout = prs.slide_layouts[1]  # 假設1是純文字版型
    elif layout_type == "title_with_content":
        slide_layout = prs.slide_layouts[1]  # 假設2是標題和內容版型
    elif layout_type == "two_column":
        slide_layout = prs.slide_layouts[3]  # 假設3是雙欄版型
    elif layout_type == "blank":
        slide_layout = prs.slide_layouts[6]  # 假設3是雙欄版型

    # 添加新頁面，並確保頁面插入到封面和封底之間
    slide = prs.slides.add_slide(slide_layout)
    prs.slides._sldIdLst.insert(len(prs.slides) - 1, prs.slides._sldIdLst[-1])  # 確保新頁面插入到封面和封底之間
    if slide.shapes.title:
        slide.shapes.title.text = title

    if content:
        if layout_type == "title_with_content":
            textbox = slide.shapes[-1]
        else:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        text_frame = textbox.text_frame
        if isinstance(content,list):
            content='\n'.join(content)

        text_frame.text = content
    return prs

# 3. 更新頁面圖像
def update_slide_image(slide, ppt_slide_outline):
    def resize_image(image, size):
        if isinstance(image,str) and os.path.exists(image) and os.path.isfile(image):
            image=Image.open(image)
        # 確保保持圖片的長寬比
        original_width, original_height = image.size
        aspect_ratio = original_width / original_height

        # 計算新的寬度和高度，保持長寬比
        if size[0] / size[1] > aspect_ratio:
            new_height = size[1]
            new_width = int(new_height * aspect_ratio)
        else:
            new_width = size[0]
            new_height = int(new_width / aspect_ratio)

        if isinstance(image,str) and os.path.exists(image):
            image=Image.open(image)
        return image.resize((new_width,new_height), Image.ANTIALIAS)

    presentation = slide.part.package.presentation
    # 根據 prompt 生成圖片（假設有一個生成圖片的函數 generate_image）
    if "background_image" in ppt_slide_outline:
        background_image = text2im({"prompt":ppt_slide_outline["background_image"],"save_folder":os.path.join(cxt.get_paradoxism_dir(),'pptx')})
        background_image=resize_image(background_image,presentation.slide_width,presentation.slide_height)
        slide.background_image = background_image
    if "image" in ppt_slide_outline:
        image = text2im({"prompt":ppt_slide_outline["image"]["prompt"],"save_folder":os.path.join(cxt.get_paradoxism_dir(),'pptx')})

        image = resize_image(image, presentation.slide_width,presentation.slide_height)
        new_width,new_height=image.size
        # 添加圖片到 slide，保持比例，並確保文字在圖片之上
        picture = slide.shapes.add_picture(image, position["x"], position["y"], width=new_width, height=new_height)
        slide.shapes._spTree.remove(picture._element)  # 將圖片移至最底層
        slide.shapes._spTree.insert(2, picture._element)  # 保證圖片在其他元素之下
        print('圖片生成完成{0}'.format(ppt_slide_outline["image"]["prompt"]))
    return slide


# 4. 插入圖表
def insert_chart(slide, chart_type, data):
    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    chart_data.add_series('Series 1', data["values"])

    if chart_type == "bar":
        chart_type_enum = XL_CHART_TYPE.BAR_CLUSTERED
    elif chart_type == "line":
        chart_type_enum = XL_CHART_TYPE.LINE_MARKERS
    elif chart_type == "pie":
        chart_type_enum = XL_CHART_TYPE.PIE

    # 確保保持圖表的長寬比
    max_width = Inches(6)
    max_height = Inches(4)
    aspect_ratio = max_width / max_height

    # 計算新的寬度和高度，保持長寬比
    if max_width / max_height > aspect_ratio:
        new_height = max_height
        new_width = int(new_height * aspect_ratio)
    else:
        new_width = max_width
        new_height = int(new_width / aspect_ratio)

    x, y = Inches(2), Inches(2)
    chart = slide.shapes.add_chart(chart_type_enum, x, y, new_width, new_height, chart_data)
    slide.shapes._spTree.remove(chart._element)  # 將圖表移至最底層
    slide.shapes._spTree.insert(2, chart._element)  # 保證圖表在其他元素之下



def generate_ppt(ppt_json):

    if not isinstance(ppt_json,str) and isinstance(ppt_json,(dict,list)):
        pass
    if isinstance(ppt_json,str) and os.path.isfile(ppt_json):
        ppt_json=json.load(open(ppt_json,'r'))
    elif isinstance(ppt_json,str):
        ppt_json=eval(ppt_json)
    cover=ppt_json[0]
    file_path = os.path.join(cxt.get_paradoxism_dir(),'pptx',cover['title'] + '_' + get_time_suffix() + '.pptx')
    make_dir_if_need(split_path(file_path)[0])
    prs=initialize_presentation(cover['title'],cover['subtitle'] if 'subtitle' in cover else '')

    for page in ppt_json[1:]:
        try:
            page_num=page['page_num']
            title=page['title']
            if 'content' in page:
                content=page['content']
                prs=suggest_base_slide(prs,page_num, "title_with_content", title, content)
            elif 'table' in page:
                content = page['table']
                prs = suggest_base_slide(prs, page_num, "title_with_content", title, content)
            else:
                prs = suggest_base_slide(prs, page_num, "blank", title, content)
            print(f'第{page_num} :{title} 基礎投影片生成完成', flush=True)
        except:
            PrintException()
            prs.save(file_path)


    for page,slide in zip(ppt_json, prs.slides):
        try:
            page_num = page['page_num']
            if 'image' in page or "background_image" in page:
                update_slide_image(slide, page)
                print(f'第{page_num} :投影片圖像生成更新', flush=True)
        except :
            PrintException()
            prs.save(file_path)
    prs.save(file_path)




